import fitz
import pdfplumber
import io
import re
import os
import docx
import openpyxl
from pptx import Presentation

def extract_text(file_bytes: bytes, filename: str) -> dict:
    result = {
        'filename': filename,
        'text': '',
        'page_count': 0,
        'char_count': 0,
        'word_count': 0,
        'method': '',
        'error': None,
        'metadata': {}
    }
    
    ext = os.path.splitext(filename)[1].lower()
    
    try:
        if ext == '.pdf':
            result = _extract_pdf(file_bytes, result)
        elif ext in ['.docx', '.doc']:
            result = _extract_docx(file_bytes, result)
        elif ext in ['.xlsx', '.xls']:
            result = _extract_xlsx(file_bytes, result)
        elif ext in ['.pptx', '.ppt']:
            result = _extract_pptx(file_bytes, result)
        else:
            result['error'] = f"Unsupported file type: {ext}"
    except Exception as e:
        result['error'] = f"Extraction failed: {str(e)}"
        
    return result

def _extract_pdf(file_bytes, result):
    try:
        doc = fitz.open(stream=file_bytes, filetype='pdf')
        result['page_count'] = len(doc)
        result['metadata'] = extract_metadata(file_bytes, '.pdf')
        text_parts = [page.get_text() for page in doc]
        full_text = '\n\n'.join(text_parts)
        
        if len(full_text.strip()) > 100:
            result['method'] = 'pymupdf'
            return _finalize_text(full_text, result)
            
    except Exception as e_fitz:
        pass
        
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            result['page_count'] = len(pdf.pages)
            text_parts = [page.extract_text() or '' for page in pdf.pages]
            full_text = '\n\n'.join(text_parts)
            result['method'] = 'pdfplumber'
            return _finalize_text(full_text, result)
    except Exception as e_plumber:
        result['error'] = f"PDF extraction failed"
    return result

def _extract_docx(file_bytes, result):
    doc = docx.Document(io.BytesIO(file_bytes))
    result['page_count'] = 1 # Approximation
    result['metadata'] = {
        'author': doc.core_properties.author or '',
        'title': doc.core_properties.title or '',
        'created': str(doc.core_properties.created) or ''
    }
    text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = '\n\n'.join(text_parts)
    result['method'] = 'python-docx'
    return _finalize_text(full_text, result)

def _extract_xlsx(file_bytes, result):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    result['page_count'] = len(wb.sheetnames)
    result['metadata'] = {
        'author': wb.properties.creator or '',
        'title': wb.properties.title or ''
    }
    text_parts = []
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            row_text = ' '.join([str(cell) for cell in row if cell is not None])
            if row_text.strip():
                text_parts.append(row_text)
    full_text = '\n'.join(text_parts)
    result['method'] = 'openpyxl'
    return _finalize_text(full_text, result)

def _extract_pptx(file_bytes, result):
    prs = Presentation(io.BytesIO(file_bytes))
    result['page_count'] = len(prs.slides)
    result['metadata'] = {
        'author': prs.core_properties.author or '',
        'title': prs.core_properties.title or ''
    }
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
    full_text = '\n\n'.join(text_parts)
    result['method'] = 'python-pptx'
    return _finalize_text(full_text, result)

def extract_metadata(file_bytes: bytes, ext: str) -> dict:
    if ext != '.pdf': return {}
    try:
        doc = fitz.open(stream=file_bytes, filetype='pdf')
        meta = doc.metadata
        return {
            'author': meta.get('author', ''),
            'title': meta.get('title', ''),
            'subject': meta.get('subject', ''),
            'creator': meta.get('creator', ''),
            'creationDate': meta.get('creationDate', ''),
            'modDate': meta.get('modDate', ''),
            'keywords': meta.get('keywords', '')
        }
    except Exception:
        return {}

def _finalize_text(text: str, result: dict) -> dict:
    if not text:
        text = ""
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    cleaned_text = text.strip()
    result['text'] = cleaned_text
    result['char_count'] = len(cleaned_text)
    result['word_count'] = len(cleaned_text.split())
    return result
